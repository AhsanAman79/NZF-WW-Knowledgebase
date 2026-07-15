"""Text extraction from many document formats.

Full text is extracted where feasible. Legacy binary formats without a
light extractor (.doc, .ppt, iWork) are still accepted and stored; they are
made findable via their metadata (handled in the ingestion pipeline).
"""
from __future__ import annotations

import io
import re
import zipfile
from pathlib import Path

# Accepted for upload (stored + indexed). Extraction quality varies by format.
SUPPORTED_EXTENSIONS = {
    # Text / web
    ".txt", ".md", ".markdown", ".csv", ".html", ".htm", ".vtt", ".rtf",
    # Word processing
    ".doc", ".docx", ".odt", ".pages",
    # Spreadsheets
    ".xls", ".xlsx", ".ods", ".numbers",
    # Presentations
    ".ppt", ".pptx", ".odp", ".key",
    # PDF
    ".pdf",
}

# Formats we accept and store but cannot full-text extract with light libraries.
# They are indexed by metadata (title/filename/description) in the pipeline.
_METADATA_ONLY = {".doc", ".ppt", ".pages", ".numbers", ".key"}


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    try:
        if ext in _METADATA_ONLY:
            return ""  # accepted; indexed via metadata
        if ext == ".pdf":
            return _from_pdf(data)
        if ext == ".docx":
            return _from_docx(data)
        if ext == ".xlsx":
            return _from_xlsx(data)
        if ext == ".pptx":
            return _from_pptx(data)
        if ext == ".xls":
            return _from_xls(data)
        if ext in {".odt", ".ods", ".odp"}:
            return _from_odf(data)
        if ext == ".rtf":
            return _from_rtf(data)
        if ext in {".html", ".htm"}:
            return _from_html(data)
        if ext == ".vtt":
            return _from_vtt(data)
        if ext in {".txt", ".md", ".markdown", ".csv"}:
            return _decode(data)
    except Exception as exc:  # noqa: BLE001 - extraction is best-effort
        return f"[extraction error for {filename}: {exc}]"
    return ""


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return _clean("\n".join((page.extract_text() or "") for page in reader.pages))


def _from_docx(data: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return _clean("\n".join(parts))


def _from_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return _clean("\n".join(parts))


def _from_xls(data: bytes) -> str:
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"# Sheet: {sheet.name}")
        for r in range(sheet.nrows):
            cells = [str(c) for c in sheet.row_values(r) if str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _clean("\n".join(parts))


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
    return _clean("\n".join(parts))


def _from_odf(data: bytes) -> str:
    """OpenDocument (.odt/.ods/.odp): text lives in content.xml inside the zip."""
    from bs4 import BeautifulSoup

    with zipfile.ZipFile(io.BytesIO(data)) as z:
        with z.open("content.xml") as f:
            xml = f.read()
    soup = BeautifulSoup(xml, "html.parser")
    return _clean(soup.get_text(separator="\n"))


def _from_rtf(data: bytes) -> str:
    """Very light RTF text extraction: strip control words and groups."""
    text = _decode(data)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)      # hex escapes
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)      # control words
    text = text.replace("{", " ").replace("}", " ").replace("\\", " ")
    return _clean(text)


def _from_html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(data), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return _clean(soup.get_text(separator="\n"))


def _from_vtt(data: bytes) -> str:
    text = _decode(data)
    lines: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit():
            continue
        line = re.sub(r"<[^>]+>", "", line).strip()
        if line:
            lines.append(line)
    deduped: list[str] = []
    for line in lines:
        if not deduped or deduped[-1] != line:
            deduped.append(line)
    return _clean("\n".join(deduped))


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
